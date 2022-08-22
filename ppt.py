from pptx import Presentation
from PIL import Image
import io, glob, configparser


class Config:
	FILENAME = 'CONFIG.txt'
	HEADER = 'DEFAULT'
	def __init__(self):
		self.config = configparser.ConfigParser()
		self.config.read(self.FILENAME)

	def save(self):
		try:
			with open('CONFIG.txt', 'w') as configfile:
				self.config.write(configfile)
		except:
			print('CONFIG ERROR')

	def append(self, **kwargs):
		for k,v in kwargs.items():
			self.config[self.HEADER][k] = str(v)

	def get(self, key):
		return self.config[self.HEADER].get(key, None)


def fit_image(image, height=500, width=1000, copy=True):
	if height==1:height=image.height
	if width==1:width=image.width
	wr, hr = width/image.width, height /image.height
	ratio = min([wr, hr])
	output_image = image if not copy else image.copy()
	output_image = output_image.resize( (int(image.width*ratio) or 1, int(image.height*ratio) or 1), Image.Resampling.LANCZOS)
	return output_image

config = Config()

prs = Presentation('Template.pptx')
# title
slide = prs.slides[0]#.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = config.get('title')
subtitle.text = "objectives:\n"+'\n'.join([f'{i+1}) {j}' for i,j in enumerate(filter(bool,[config.get('obj1'), config.get('obj2'), config.get('obj3')]))])
for i, im in enumerate(sorted(glob.glob('*.png'))):
	slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(slide_layout)
	title,p2 = slide.placeholders
	title.text = f'Example {i+1}'
	with io.BytesIO() as output:
		image = Image.open(im)
		image = fit_image(image)
		frame=Image.new(size=(1000,500),color='white', mode='RGBA')  
		frame.paste(image, (0, 250-image.height//2)) 
		print(type(output),5)
		frame.save(output, format="PNG") 
		picture = p2.insert_picture(output)

	picture.crop_top = 0
	picture.crop_left = 0
	picture.crop_bottom = 0
	picture.crop_right = 0
	
slide_layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(slide_layout)

for p in slide.placeholders:
	p.text = f'School book page {config.get("schoolbook_hw")}\nbooklet page {config.get("booklet_hw")}'

prs.save(fr'{config.get("filename")}.pptx')
